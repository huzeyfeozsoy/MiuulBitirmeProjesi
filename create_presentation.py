from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Sunum oluştur
prs = Presentation()

# Temel stiller (Koyu tema / Yeşil vurgu)
def style_title(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 255, 135)

def style_body(shape):
    for paragraph in shape.text_frame.paragraphs:
        paragraph.space_after = Pt(14)
        for run in paragraph.runs:
            run.font.size = Pt(20)

# Slayt 1: Kapak & Problem Tanımı
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "Slayt 1: Problem Tanımı (İş Problemi Nedir?)"
body.text = "• Futbol Kulüpleri İçin Risk: Yanlış transfer kararlarının getirdiği milyonlarca euroluk ekonomik ve sportif zararlar.\n" \
            "• Mevcut Durum: Gözleme ve menajer yönlendirmesine dayalı eski usul scouting yöntemleri yetersiz.\n" \
            "• Çözümümüz (SmartClub 360): Hem piyasa değeri tahmini hem de oyuncunun 'Baskı altındaki' gerçek bitiricilik yeteneğini (Contextual xG) hesaplayan, veri odaklı uçtan uca scouting portalı."
style_title(title)
style_body(body)

# Slayt 2: Dataset Özeti
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "Slayt 2: Dataset Özeti"
body.text = "• Veri Kaynakları: Transfermarkt (Piyasa) ve Understat (Performans) (2025/2026 Sezonu)\n" \
            "• Boyutlar: 839 Oyuncu (Streamlit DB), ~35 Feature (Dakika, Gol, Asist, xG, xA, Şut vb.)\n" \
            "• Hedef Değişken (Target 1): 'market_value' (Piyasa Değeri - Euro) Tahmini\n" \
            "• Hedef Değişken (Target 2): Contextual xG modeli için 'Gol Olma İhtimali' (Sınıflandırma)"
style_title(title)
style_body(body)

# Slayt 3: EDA Bulguları
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "Slayt 3: EDA Bulguları (En Çarpıcı 3 Insight)"
body.text = "1. Yaş ve Değerin Ters Orantısı: 25 yaşından sonra piyasa değerinde dramatik bir düşüş eğilimi gözlemleniyor.\n" \
            "2. Dakika Kısıtı: Minimum 400 dakika oynamayan oyuncuların verileri aşırı gürültülü. Optimizasyon için 400 dakika barajı koyuldu.\n" \
            "3. Bağlamsal (Contextual) Fark: Baskı altındaki şutların gol olma oranı, baskısız şutlara göre %40 daha düşük. Klasik xG bu durumu yakalayamıyor."
style_title(title)
style_body(body)

# Slayt 4: Feature Engineering Showcase
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "Slayt 4: Feature Engineering Showcase"
body.text = "• Yaş Potansiyeli Skoru (Age Potential Score): (Gol + Asist / Dakika * 90) / Yaş hesaplaması ile genç yeteneklerin tespit edilmesi.\n" \
            "• 90 Dakika Metrikleri: xG_xA_per90, shots_per90, key_passes_per90 ile oyuncu süreleri normalize edildi.\n" \
            "• Contextual Variables: Şut çekilen andaki 'Oyun Skoru' ve 'Rakip Baskı Seviyesi' (PPDA gibi) modele eklendi.\n" \
            "• Veri Kurtarma (Data Enrichment): 113 oyuncunun API'ler arası isim farkı (Ör: Rayan Cherki vs Mathis Cherki) NLP mantığıyla çözüldü."
style_title(title)
style_body(body)

# Slayt 5: Model Karşılaştırma Tablosu
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "Slayt 5: Model Karşılaştırma Tablosu"
body.text = "• Denenen Modeller: Linear Regression, Random Forest, XGBoost, LightGBM.\n" \
            "• Piyasa Değeri Regresyonu İçin: XGBoost ve LightGBM en iyi sonuçları verdi.\n" \
            "• Aşırı Öğrenme / Halüsinasyon Kontrolü: Regresyon modeli bazen uçuk tahminler (örn: %9500 fark) verdiğinde, iş dünyasına uyarlama (Shrinkage: %25 Model + %75 Gerçek Değer) stratejisi kullanıldı.\n" \
            "• xG Sınıflandırma (Contextual): Logistic Regression ve XGBoost Classifier ile o an gol olma ihtimali karşılaştırıldı."
style_title(title)
style_body(body)

# Slayt 6: Final Model Metrikleri & Feature Importance
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "Slayt 6: Final Model Metrikleri & Feature Importance"
body.text = "• Final Model (Piyasa Değeri): LightGBM Regressor (CV R²: ~0.61)\n" \
            "• En Önemli 4 Değişken (Feature Importance):\n" \
            "   1. Yaş (dob_age)\n" \
            "   2. 90 Dakika Başına Gol Katkısı (goal_contribution_per90)\n" \
            "   3. Oynanan Dakika (minutes)\n" \
            "   4. 90 Dakika Başına xG+xA\n" \
            "• Model Yorumlanabilirliği: SHAP analizi ile oyuncuların karar anlarında hangi özelliklerinin öne çıktığı şeffaflaştırıldı."
style_title(title)
style_body(body)

# Slayt 7: Business Impact / ROI Analizi
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "Slayt 7: Business Impact / ROI Analizi"
body.text = "• Fırsat Transferlerinin Keşfi: Gerçek değeri €10M olan ancak modelin €15M potansiyel gördüğü (Underpriced) oyuncuların erken transferi.\n" \
            "• Risk Yönetimi: Çok para istenen ama performansı düşen (Overpriced) oyunculardan kaçınarak milyonlarca Euro tasarruf.\n" \
            "• Zihinsel Scouting: Contextual xG sayesinde sadece yetenekli değil, 'Zor anlarda' iş yapan karakterli oyuncuların tespiti ile takım ROI'sinin maksimize edilmesi."
style_title(title)
style_body(body)

# Slayt 8: Streamlit Demo
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "Slayt 8: Streamlit Demo (Canlı Göster!)"
body.text = "• Futbol Temalı Modern Arayüz: Koyu yeşil çim dokusu, neon detaylar.\n" \
            "• Gerçek Zamanlı API: Fantasy Premier League API üzerinden 20 kulübün logoları ve 839 oyuncunun yüzleri dinamik çekiliyor.\n" \
            "• 'Kuş Bakışı Saha' Simülasyonu: Seçilen mevkiye göre dinamik oyuncu filtreleme ve profil kartları.\n" \
            "• Modüler Yapı: Taha'nın Piyasa Değeri Modeli ile Huzeyfe'nin Contextual xG modeli tek portalda uyum içinde çalışıyor."
style_title(title)
style_body(body)

# Slayt 9: Sonuç & Öğrenilenler
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "Slayt 9: Sonuç & Öğrenilenler"
body.text = "• Veri Kalitesi Modeli Yener: Çeşitli kaynaklardan gelen verilerde isim uyumsuzluklarının (Transfermarkt vs Understat) modelden 113 oyuncu eksilttiğini görmek büyük bir dersti.\n" \
            "• Saf Matematik Her Şey Değildir: Modelin teknik skoru (R²) dışında, sonuçları mantıklı sınırlara çekmek (Shrinkage) veri biliminin iş dünyasına satışını kolaylaştırır.\n" \
            "• Ekip İçi Entegrasyon: İki farklı analitik modelin birleşerek tek bir güçlü ürüne (Scouting Portal) dönüşmesinin gücü."
style_title(title)
style_body(body)

prs.save("Miuul_Bitirme_Projesi_Sunumu.pptx")
print("Sunum başarıyla oluşturuldu: Miuul_Bitirme_Projesi_Sunumu.pptx")
